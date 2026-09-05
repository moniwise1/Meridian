"""
Presentation Generator (BUILD SPEC section 24). Builds a management deck
from the same already-computed result snapshot the report uses - every
chart/number on a slide corresponds to an actual analytical result, never
fabricated to "fill" a slide.
"""
import os
import uuid
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from app.config import settings


def _add_branding(prs, slide) -> None:
    """A small "Made by Meridian" footer on every slide - python-pptx has
    no per-slide auto-callback the way fpdf2 does, so this is called
    explicitly after each slide is built. Bottom-right, out of the way of
    the actual content (title/bullets/table all sit above this band).
    Purely cosmetic; never touches the slide's real content."""
    box = slide.shapes.add_textbox(
        prs.slide_width - Inches(3.2), prs.slide_height - Inches(0.4), Inches(3.0), Inches(0.3),
    )
    tf = box.text_frame
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "Made by Meridian"
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)


def _add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    _add_branding(prs, slide)
    return slide


def _add_bullets_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = b
        p.font.size = Pt(16)
    _add_branding(prs, slide)
    return slide


def _add_table_slide(prs, title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    rows_n = min(len(rows), 12) + 1
    cols_n = len(headers)
    table_shape = slide.shapes.add_table(rows_n, cols_n, Inches(0.5), Inches(1.5), Inches(9), Inches(0.4 * rows_n))
    table = table_shape.table
    for c, h in enumerate(headers):
        table.cell(0, c).text = str(h)
    for r, row in enumerate(rows[:12], start=1):
        for c, val in enumerate(row):
            table.cell(r, c).text = str(val)
    _add_branding(prs, slide)
    return slide


def generate_presentation_pptx(title: str, question: str, insight: dict, metrics: dict,
                                 by_group: list[dict] | None, data_quality: dict,
                                 anomalies: list[dict], query_id: str) -> str:
    os.makedirs(settings.artifacts_dir, exist_ok=True)
    path = os.path.join(settings.artifacts_dir, f"presentation-{uuid.uuid4().hex[:8]}.pptx")

    prs = Presentation()

    _add_title_slide(prs, title, f"{question}\nQuery ID: {query_id}")

    if "error" not in insight:
        _add_bullets_slide(prs, "Executive summary", [
            insight.get("what", ""),
            f"Where: {insight.get('where', '')}",
            f"When: {insight.get('when', '')}",
        ])
        _add_bullets_slide(prs, "Key findings", [
            insight.get("contributors", ""),
            f"Confidence: {insight.get('confidence', '')} — {insight.get('confidence_explanation', '')}",
            f"Next question: {insight.get('next_question', '')}",
        ])

    if by_group:
        _add_table_slide(
            prs, "Breakdown", ["Group", "Total"],
            [[row["group"], f"{row['total']:,.2f}"] for row in by_group],
        )

    if anomalies:
        _add_bullets_slide(prs, "Risks & anomalies", [
            f"{a['what']} — {a['magnitude']} [{a['confidence']} confidence]" for a in anomalies[:5]
        ])

    _add_bullets_slide(prs, "Data quality & methodology", [
        f"Rows analysed: {data_quality.get('row_count', 0)}",
        f"Completeness: {data_quality.get('completeness_pct', 100)}%",
        *[f"Note: {n}" for n in data_quality.get("notes", [])],
        "All figures computed deterministically; the AI model interprets results, it does not calculate them.",
    ])

    prs.save(path)
    return path
