"""
Document Intelligence — extracts plain text from an uploaded PDF, DOCX,
PPTX, or XLSX so it can be referenced alongside a database analysis (BUILD
SPEC section 19). This is text/table extraction, not comprehension: it
turns a file into text the insight-explanation LLM step can read, exactly
the way it already reads computed metrics — nothing here interprets,
summarizes, or validates the document's content.

Security note this module exists specifically to keep in view: extracted
document text is the first genuinely externally-authored content this app
ever hands to an LLM. Row values and schema field names are trusted enough
(they come from a database the tenant connected and authorized), but a
document a user uploads could contain anything, including text specifically
crafted to look like instructions ("ignore your previous instructions and
...", "SYSTEM:", etc.). Per app/agents/planner.py's prompt-injection
defence, this text must always be handed to the LLM as a labelled,
untrusted DATA payload the model is told to reference, never blended into
instruction text — see how app/agents/insight_agent.py passes it. This
module's job stops at extraction; it does not decide how the text is used
downstream.

OCR fallback for scanned/image-only PDF pages: pypdf only reads an
embedded text layer, so a genuinely scanned page (a photo/scan with no
text layer at all) used to extract to nothing, silently. Now, any page
whose native extraction comes back empty is rendered to an image
(PyMuPDF - no external renderer binary needed, unlike poppler) and
run through Tesseract (pytesseract - a real external OCR engine binary,
NOT pip-installable on its own; see docs/OCR.md for the one system-level
install step this needs, already wired into backend/Dockerfile for
production). Deliberately per-PAGE, not per-document: a mixed PDF (some
real text pages, some scanned pages - a common real-world shape, e.g. a
native report with a scanned signature page appended) gets native
extraction for the pages that have it and OCR only for the pages that
need it, rather than an all-or-nothing choice. Bounded by
MAX_OCR_PAGES_PER_DOCUMENT since OCR is genuinely CPU-expensive (roughly
1-3s/page at the DPI used here) unlike the near-instant native path, and
this upload endpoint is still synchronous - no background job queue
exists in this app to hand slow work off to, so an unbounded scanned PDF
could otherwise stall the request for minutes. Fails open, not closed, if
Tesseract isn't installed on the machine at all (TesseractNotFoundError):
falls back to the pre-OCR behavior (empty text for that page, surfaced
honestly, never pretending it worked) rather than crashing the upload -
same "a missing optional capability degrades, it doesn't break the app"
pattern already used for Redis and the Anthropic client elsewhere in this
app.

Still NOT built: real PDF table structure (text extraction, OCR'd or
native, flattens tables into reading-order text, which reads poorly for
anything but simple layouts - a known, unfixed limitation, not a bug).
"""
import io
import logging
from dataclasses import dataclass

import pypdf
import docx
import openpyxl
import pptx
import pytesseract
import pymupdf

from app.config import settings

logger = logging.getLogger("meridian.ocr")

if settings.tesseract_cmd:
    # Only needed where Tesseract isn't already resolvable on PATH - e.g.
    # local Windows dev, where its installer doesn't always add itself to
    # PATH for an already-open shell. The Linux/Docker production target
    # (backend/Dockerfile installs tesseract-ocr via apt) needs no
    # override at all; PATH resolution just works there.
    pytesseract.pytesseract.tesseract_cmd = settings.tesseract_cmd

MAX_EXTRACTED_CHARS = 50_000  # bounds LLM context cost the same way row limits bound query cost
MAX_XLSX_ROWS_PER_SHEET = 200
MAX_XLSX_SHEETS = 10
MAX_PPTX_SLIDES = 200
MAX_OCR_PAGES_PER_DOCUMENT = 15
OCR_RENDER_DPI = 200  # balance of accuracy vs. render+recognition time
# A page's native extraction shorter than this is treated as "probably
# scanned, not just a sparse page" and gets OCR'd - a real PDF page with
# only a few words of genuine text is rare enough that this heuristic
# costs little precision while catching the common case (an image-only
# page returns "" or a handful of stray characters from decorative
# elements, never a real sentence).
NATIVE_TEXT_MIN_CHARS = 20

SUPPORTED_EXTENSIONS = {".pdf": "pdf", ".docx": "docx", ".xlsx": "xlsx", ".pptx": "pptx"}


class UnsupportedDocumentType(Exception):
    pass


class DocumentTooLarge(Exception):
    pass


@dataclass
class ExtractionResult:
    text: str
    truncated: bool
    # Meaning depends on kind: page count for PDF, paragraph count for
    # DOCX, sheet count for XLSX — informational only, shown in the UI.
    source_unit_count: int
    # How many pages' text came from OCR rather than a native text layer
    # (PDF only; always 0 for every other kind). Surfaced to the caller so
    # the UI can show "(N page(s) OCR'd)" - OCR'd text is real but lower-
    # confidence than a native text layer (misreads happen), worth
    # flagging rather than presenting identically to a clean extraction.
    ocr_pages_used: int = 0


def _truncate(text: str) -> tuple[str, bool]:
    text = text.strip()
    if len(text) <= MAX_EXTRACTED_CHARS:
        return text, False
    return text[:MAX_EXTRACTED_CHARS], True


def _ocr_page(pdf_doc, page_index: int) -> str:
    """Renders one page to an image and runs Tesseract on it. Returns ""
    (not an exception) on any OCR-specific failure - a page that can't be
    OCR'd degrades to "no text from this page", the same honest-empty
    result a scanned page without this fallback at all would have
    produced, never a crash that takes down the whole upload over one
    bad page."""
    try:
        page = pdf_doc[page_index]
        pix = page.get_pixmap(dpi=OCR_RENDER_DPI)
        image_bytes = pix.tobytes("png")
        from PIL import Image
        image = Image.open(io.BytesIO(image_bytes))
        return pytesseract.image_to_string(image).strip()
    except pytesseract.TesseractNotFoundError:
        # Not installed on this machine at all - fail open for the WHOLE
        # document, not just this page (every subsequent OCR attempt
        # would hit the identical error), by re-raising a marker the
        # caller checks for once and stops trying further pages.
        raise
    except Exception as e:
        logger.warning("OCR failed for page %d, treating as empty: %s", page_index, e)
        return ""


def extract_pdf(file_bytes: bytes) -> ExtractionResult:
    reader = pypdf.PdfReader(io.BytesIO(file_bytes))
    pages_text = [page.extract_text() or "" for page in reader.pages]

    ocr_pages_used = 0
    candidates = [i for i, t in enumerate(pages_text) if len(t.strip()) < NATIVE_TEXT_MIN_CHARS]
    if candidates:
        tesseract_available = True
        pdf_doc = pymupdf.open(stream=file_bytes, filetype="pdf")
        try:
            for i in candidates[:MAX_OCR_PAGES_PER_DOCUMENT]:
                if not tesseract_available:
                    break
                try:
                    ocr_text = _ocr_page(pdf_doc, i)
                except pytesseract.TesseractNotFoundError:
                    logger.warning("Tesseract is not installed on this machine - OCR fallback skipped "
                                    "for the rest of this document (and every document until it is).")
                    tesseract_available = False
                    continue
                if ocr_text:
                    pages_text[i] = ocr_text
                    ocr_pages_used += 1
        finally:
            pdf_doc.close()

    text, truncated = _truncate("\n\n".join(pages_text))
    return ExtractionResult(
        text=text, truncated=truncated, source_unit_count=len(reader.pages),
        ocr_pages_used=ocr_pages_used,
    )


def extract_docx(file_bytes: bytes) -> ExtractionResult:
    document = docx.Document(io.BytesIO(file_bytes))
    parts = [p.text for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text.strip() for cell in row.cells))
    text, truncated = _truncate("\n".join(parts))
    return ExtractionResult(text=text, truncated=truncated, source_unit_count=len(document.paragraphs))


def extract_xlsx(file_bytes: bytes) -> ExtractionResult:
    workbook = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    total_sheets = len(workbook.worksheets)
    sheets_included = 0
    row_truncated = False
    parts = []
    for sheet in workbook.worksheets[:MAX_XLSX_SHEETS]:
        sheets_included += 1
        parts.append(f"--- Sheet: {sheet.title} ---")
        for i, row in enumerate(sheet.iter_rows(values_only=True)):
            if i >= MAX_XLSX_ROWS_PER_SHEET:
                parts.append(f"... (more rows omitted, showing first {MAX_XLSX_ROWS_PER_SHEET})")
                row_truncated = True
                break
            parts.append(" | ".join("" if v is None else str(v) for v in row))
    workbook.close()
    text, char_truncated = _truncate("\n".join(parts))
    truncated = char_truncated or row_truncated or sheets_included < total_sheets
    return ExtractionResult(text=text, truncated=truncated, source_unit_count=total_sheets)


def extract_pptx(file_bytes: bytes) -> ExtractionResult:
    presentation = pptx.Presentation(io.BytesIO(file_bytes))
    total_slides = len(presentation.slides)
    slides_included = 0
    parts = []
    for i, slide in enumerate(presentation.slides):
        if i >= MAX_PPTX_SLIDES:
            break
        slides_included += 1
        parts.append(f"--- Slide {i + 1} ---")
        for shape in slide.shapes:
            # Title/body text boxes, and any other shape with a text
            # frame (a caption, a text box someone dragged in, etc).
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs).strip()
                    if text:
                        parts.append(text)
            # Tables render the same "cell | cell" flattening as
            # extract_docx's tables, for consistency across formats.
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append(" | ".join(cell.text.strip() for cell in row.cells))
        # Speaker notes often carry real analytical content (the actual
        # narration a deck's bullet points only hint at) - included, but
        # clearly labelled so it's obvious in the extracted text which
        # part was on-slide vs. notes-only.
        if slide.has_notes_slide:
            notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes_text:
                parts.append(f"[Speaker notes] {notes_text}")
    text, char_truncated = _truncate("\n".join(parts))
    truncated = char_truncated or slides_included < total_slides
    return ExtractionResult(text=text, truncated=truncated, source_unit_count=total_slides)


def extract(filename: str, file_bytes: bytes) -> tuple[str, ExtractionResult]:
    """Dispatches on file extension. Returns (kind, ExtractionResult)."""
    ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    kind = SUPPORTED_EXTENSIONS.get(ext)
    if kind is None:
        raise UnsupportedDocumentType(
            f"Unsupported file type '{ext or filename}'. Supported: "
            f"{', '.join(sorted(SUPPORTED_EXTENSIONS))}.",
        )
    if kind == "pdf":
        return kind, extract_pdf(file_bytes)
    if kind == "docx":
        return kind, extract_docx(file_bytes)
    if kind == "pptx":
        return kind, extract_pptx(file_bytes)
    return kind, extract_xlsx(file_bytes)
